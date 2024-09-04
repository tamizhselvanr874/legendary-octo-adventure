import streamlit as st
import requests
import json
import os
from pathlib import Path
import fitz  # PyMuPDF
from langchain_openai import AzureChatOpenAI
from tenacity import retry, wait_random_exponential, stop_after_attempt
from docx import Document
from docx.shared import Inches
import io
import re
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Configuration
GRAPH_TENANT_ID = "4d4343c6-067a-4794-91f3-5cb10073e5b4"
GRAPH_CLIENT_ID = "5ace14db-3235-4cd2-acfd-dd5ef19d6ea1"
GRAPH_CLIENT_SECRET = "HRk8Q~7G6EH3.yhDC3rB5wLAyAixQMnQNWNyUdsW"
PDF_SITE_ID = "marketingai.sharepoint.com,b82dbaac-09cc-4539-ad08-e4ca926796e8,7b756d20-3463-44b7-95ca-5873f8c3f517"
FUNCTION_URL = "https://doc2pdf.azurewebsites.net"

# Azure OpenAI API details
azure_endpoint = 'https://chat-gpt-a1.openai.azure.com/'
azure_deployment_name = 'DanielChatGPT16k'
azure_api_key = 'c09f91126e51468d88f57cb83a63ee36'
azure_api_version = '2024-05-01-preview'

# Initialize Azure OpenAI LLM
llm = AzureChatOpenAI(
    openai_api_key=azure_api_key,
    api_version=azure_api_version,
    azure_endpoint=azure_endpoint,
    model="gpt-4",
    azure_deployment=azure_deployment_name,
    temperature=0.5
)

def get_oauth2_token():
    url = f"https://login.microsoftonline.com/{GRAPH_TENANT_ID}/oauth2/v2.0/token"
    headers = {'Content-Type': 'application/x-www-form-urlencoded'}
    data = {
        'grant_type': 'client_credentials',
        'client_id': GRAPH_CLIENT_ID,
        'client_secret': GRAPH_CLIENT_SECRET,
        'scope': 'https://graph.microsoft.com/.default'
    }
    response = requests.post(url, headers=headers, data=data)
    if response.status_code == 200:
        return response.json().get('access_token')
    else:
        st.error(f"Failed to obtain OAuth2 token: {response.content}")
        return None

def upload_file_to_sharepoint(token, file):
    upload_url = f"https://graph.microsoft.com/v1.0/sites/{PDF_SITE_ID}/drive/root:/{file.name}:/content"
    headers = {
        'Authorization': f'Bearer {token}',
        'Content-Type': file.type
    }
    response = requests.put(upload_url, headers=headers, data=file.getvalue())
    
    if response.status_code in [200, 201]:
        return response.json().get('id')
    else:
        st.error(f"Failed to upload file to SharePoint: {response.content}")
        return None

def convert_file_to_pdf(token, file_id):
    convert_url = f"https://graph.microsoft.com/v1.0/sites/{PDF_SITE_ID}/drive/items/{file_id}/content?format=pdf"
    headers = {
        'Authorization': f'Bearer {token}',
        'Content-Type': 'application/json'
    }
    response = requests.get(convert_url, headers=headers)
    if response.status_code == 200:
        return response.content
    else:
        st.error(f"Failed to convert file to PDF: {response.content}")
        return None

def delete_file_from_sharepoint(token, file_id):
    delete_url = f"https://graph.microsoft.com/v1.0/sites/{PDF_SITE_ID}/drive/items/{file_id}"
    headers = {'Authorization': f'Bearer {token}'}
    response = requests.delete(delete_url, headers=headers)
    if response.status_code == 204:
        return True
    else:
        st.error(f"Failed to delete file from SharePoint: {response.content}")
        return False

def read_prompt(prompt_path: str):
    with open(prompt_path, "r") as f:
        return f.read()

def extract_text_from_pdf(pdf_path: str):
    doc = fitz.open(pdf_path)
    text = ""
    for page_num, page in enumerate(doc):
        page_text = page.get_text()
        print(f"Extracting text from Page {page_num + 1}: {page_text[:100]}...")  # Print the first 100 characters for debugging
        text += f"Page {page_num + 1}\n" + page_text
    return text

def extract_images_from_pdf(pdf_path: str, output_folder: str):
    doc = fitz.open(pdf_path)
    images = []
    
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    for page_number in range(len(doc)):
        page = doc[page_number]
        image_list = page.get_images(full=True)
        
        for image_index, img in enumerate(image_list, start=1):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            image_path = os.path.join(output_folder, f"image_{page_number + 1}_{image_index}.{image_ext}")
            
            with open(image_path, "wb") as img_file:
                img_file.write(image_bytes)
            
            # Simulate meaningful and relevant image titles and descriptions
            image_title = f"Figure {page_number + 1}.{image_index}: Description of key feature"
            image_description = f"This figure illustrates key aspect {image_index} found on page {page_number + 1}."

            images.append({
                "page_number": page_number + 1,
                "title": image_title,
                "description": image_description,
                "image_url": image_path
            })
    
    return images

@retry(wait=wait_random_exponential(min=1, max=120), stop=stop_after_attempt(10))
def completion_with_backoff(prompt: str, content: str):
    try:
        response = llm(
            messages=[
                {"role": "system", "content": prompt},
                {"role": "user", "content": content}
            ]
        )
        return response
    except Exception as e:
        print(f"Error calling Azure OpenAI API: {e}")
        return {}

def extract_metadata(content: str):
    prompt = """Patent Document Analysis Prompt:

    You are an expert tasked with analyzing patent documents. Please thoroughly review the provided patent document and extract the following key information from each page:

    1. Page Number: Identify the page number of the document.
    2. Page Title: Extract the title or heading of the page.
    3. Page Content: Extract the exact main content and context of the page as it appears in the document, without any summarization.
    4. Image: Identify and extract any images present on the page along with relevant metadata. This includes the image title, description, and any other pertinent information.

    Guidelines:
    - Ensure that all extracted information is factual, accurate, and directly derived from the document.
    - For the "Page Title" and "Image" sections, provide concise and descriptive information.
    - The information should be self-contained, meaning that each extracted piece should make sense independently of the rest of the document.
    - If any of the required properties are not present on a page, leave those fields empty rather than making assumptions.
    - For images, include detailed metadata such as:
      - Image Title: The title or caption associated with the image.
      - Image Description: A brief description of the image’s content and purpose.
      - Additional Metadata: Any other relevant details, such as image source or reference numbers.

    Response Format:
    Answer in JSON format. Each page should be represented as an object with the following keys:

    - "PageNumber": The number of the page.
    - "PageTitle": The title of the page as it appears in the document.
    - "PageContent": The exact main content and context of the page as it appears in the document.
    - "Images": A list of objects containing:
      - "ImageTitle": The title of the image.
      - "ImageDescription": A description of the image.
      - "AdditionalMetadata": Any other relevant image metadata.

    Here is an example JSON format for your response:

    ```json
    [
      {
        "PageNumber": 1,
        "PageTitle": "Title of the Page",
        "PageContent": "Exact content of the page without any summarization.",
        "Images": [
          {
            "ImageTitle": "Title of the Image",
            "ImageDescription": "Description of the image.",
            "AdditionalMetadata": "Other relevant details about the image."
          }
        ]
      },
      ...
    ]
    ```"""


    response = completion_with_backoff(prompt, content)
    if not response:
        print("Empty response from the model")
        return []

    if hasattr(response, 'content'):
        response_content = response.content
        response_content = re.sub(r'```json\s*', '', response_content)
        response_content = re.sub(r'\s*```', '', response_content)

        print(f"Raw JSON response: {response_content}")

        try:
            metadata = json.loads(response_content)
        except json.JSONDecodeError as e:
            print(f"Error decoding JSON: {e}")
            metadata = []

    return metadata

def extract_tables_and_flow_diagrams_from_pptx(pptx_path):
    presentation = Presentation(pptx_path)
    tables_and_diagrams = []

    for slide_index, slide in enumerate(presentation.slides):
        slide_data = {
            "page_number": slide_index + 1,
            "tables": [],
            "flow_diagrams": []
        }
        
        for shape in slide.shapes:
            # Extract tables
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                slide_data["tables"].append(table_data)
            
            # Extract flow diagrams (shapes)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                diagram_data = []
                for s in shape.shapes:
                    if s.has_text_frame:
                        diagram_data.append(s.text)
                slide_data["flow_diagrams"].append(diagram_data)
        
        if slide_data["tables"] or slide_data["flow_diagrams"]:
            tables_and_diagrams.append(slide_data)
    
    return tables_and_diagrams

def create_word_file_with_structured_content(json_data, pptx_path):
    doc = Document()
    doc.add_heading('Extracted Metadata', 0)

    # Extract tables and diagrams from the PPTX file
    structured_content = extract_tables_and_flow_diagrams_from_pptx(pptx_path)

    for page in json_data:
        doc.add_heading(f"Page {page['PageNumber']}", level=1)
        doc.add_heading('Header', level=2)
        doc.add_paragraph(page['PageTitle'])
        doc.add_heading('Content', level=2)
        doc.add_paragraph(page['PageContent'])

        # Integrate structured content (tables and diagrams)
        for structured_page in structured_content:
            if structured_page["page_number"] == page['PageNumber']:
                if structured_page["tables"]:
                    doc.add_heading('Tables', level=2)
                    for table in structured_page["tables"]:
                        table_str = "\n".join(["\t".join(row) for row in table])
                        doc.add_paragraph(table_str)

                if structured_page["flow_diagrams"]:
                    doc.add_heading('Flow Diagrams', level=2)
                    for diagram in structured_page["flow_diagrams"]:
                        doc.add_paragraph("\n".join(diagram))
        
        if 'Images' in page and page['Images']:
            for image in page['Images']:
                doc.add_heading('Image', level=2)
                doc.add_paragraph(f"Title: {image['ImageTitle']}")
                doc.add_paragraph(f"Description: {image['ImageDescription']}")
                doc.add_paragraph(f"Additional Metadata: {image.get('AdditionalMetadata', 'N/A')}")
                # Add image to the document
                doc.add_picture(image['image_url'], width=Inches(4))
    
    file_stream = io.BytesIO()
    doc.save(file_stream)
    file_stream.seek(0)
    return file_stream

# Streamlit app
st.title("PPT to Structured PDF and Word Document Converter")

uploaded_file = st.file_uploader("Upload a PPT file", type=["pptx"])

if uploaded_file is not None:
    token = get_oauth2_token()
    
    if token:
        # Upload PPT to SharePoint
        file_id = upload_file_to_sharepoint(token, uploaded_file)
        
        if file_id:
            # Convert PPT to PDF
            pdf_content = convert_file_to_pdf(token, file_id)
            
            if pdf_content:
                # Save PDF locally
                pdf_path = Path("output.pdf")
                with open(pdf_path, "wb") as pdf_file:
                    pdf_file.write(pdf_content)
                
                # Extract text and images
                extracted_text = extract_text_from_pdf(str(pdf_path))
                images = extract_images_from_pdf(str(pdf_path), "extracted_images")
                
                # Extract metadata
                metadata = extract_metadata(extracted_text)

                if metadata:
                    for image in images:
                        for page in metadata:
                            if page['PageNumber'] == image['page_number']:
                                if 'Images' not in page or not page['Images']:
                                    page['Images'] = []
                                page['Images'].append({
                                    "ImageTitle": image['title'],
                                    "ImageDescription": image['description'],
                                    "image_url": image['image_url']
                                })

                    st.json(metadata)
                    
                    # Create and download Word file with structured content
                    word_file = create_word_file_with_structured_content(metadata, uploaded_file)
                    st.download_button(
                        label="Download Word file",
                        data=word_file,
                        file_name="extracted_metadata.docx",
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    )
                
                # Clean up SharePoint by deleting the uploaded file
                delete_file_from_sharepoint(token, file_id)
            else:
                st.error("PDF conversion failed.")
        else:
            st.error("File upload to SharePoint failed.")
